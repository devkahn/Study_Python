from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
for i in range(0,11):
    print("--------[%d] ------ "%(i))
    slide = prs.slides.add_slide(prs.slide_layouts[i])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))



