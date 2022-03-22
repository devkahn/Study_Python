from copyreg import pickle
from turtle import width
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()


title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

title = slide.placeholders[0]
title.text = "Hello, world!"

subtitle = slide.placeholders[1]
subtitle.text = "python-pptx was here!!!"


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)

title_shape = slide.placeholders[0]
title_shape.text = 'Adding a Bullet Slide'

body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use_TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for susequent bullets'
p.level = 2

img_path = '222.gif'

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
width = height = Inches(1)

pic = slide.shapes.add_picture(img_path, left, top, width=width, height = height)

left = Inches(3)
width = Inches(5.5)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, width = width, height = height)


title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

title_shape = slide.placeholders[0]
title_shape.text ='Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

table.cell(0,0).text = 'Foo'
table.cell(0,1).text = 'Bar'

table.cell(1,0).text = 'Baz'
table.cell(1,1).text = 'Qux'





prs.save('test.pptx')