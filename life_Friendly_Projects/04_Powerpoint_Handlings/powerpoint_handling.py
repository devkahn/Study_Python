from turtle import title
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

for i in range(0,11):
    title_slide_layout  = prs.slide_layouts[i]
    slide = prs.slides.add_slide(title_slide_layout)

prs.save('add all slides.pptx')